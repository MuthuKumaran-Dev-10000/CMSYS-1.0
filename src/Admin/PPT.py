from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]  # 0 is the layout for title slides
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "ONDC Application Development Plan"
subtitle.text = "Project Phases and Sprint Breakdown"

# Add content slide (Phase 1)
slide_layout = prs.slide_layouts[1]  # 1 is the layout for content slides
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Phase 1: AI Feature Development"

content = slide.shapes.placeholders[1]
content.text = (
    "1. Initial development of AI suggestion features (Jan - Mar)\n"
    "2. Implementation of statistical data (Feb - Apr)\n"
    "3. UI development for trend display (Mar - Jun)\n"
    "4. Testing and refinement of AI models (Jun - Aug)\n"
    "5. Notification system for real-time trends (May - Jun)\n"
    "6. Full AI-based recommendation system launch (Jul - Sep)"
)

# Add content slide (Phase 2)
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Phase 2: Expansion and Feedback Integration"

content = slide.shapes.placeholders[1]
content.text = (
    "1. Personalized seller dashboards (Apr - Sep)\n"
    "2. Backend optimization for real-time prediction (Jan - Mar)\n"
    "3. Customer feedback collection system (Apr - Jun)\n"
    "4. Feature upgrades for seller-customer engagement (Aug - Oct)\n"
    "5. Major milestone: Fully functional version launch (Sep)"
)

# Add content slide (Phase 3)
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Phase 3: Advanced Analytics"

content = slide.shapes.placeholders[1]
content.text = (
    "1. Development of advanced analytics (Jan - Jul)\n"
    "2. Report on AI-driven suggestions (Due: 05/20)\n"
    "3. Public launch of analytics features (Launch: 07/01)\n"
    "4. Further AI model iterations (Feb - Aug)"
)

# Add slide for Sprint Planning
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Sprint Planning Breakdown"

content = slide.shapes.placeholders[1]
content.text = (
    "Sprint 1:\n"
    "- Develop AI trend analysis engine\n"
    "- Build seller dashboard\n"
    "- Integrate statistical data analysis\n\n"
    
    "Sprint 2:\n"
    "- Improve AI model accuracy\n"
    "- Implement notification system\n\n"
    
    "Sprint 3:\n"
    "- Personalization algorithm development\n"
    "- UI enhancements\n\n"
    
    "Sprint 4:\n"
    "- Final testing and optimization\n"
    "- Prepare for public launch"
)

# Add slide for Additional Comments
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Additional Comments"

content = slide.shapes.placeholders[1]
content.text = (
    "1. Real-time trend suggestions will give sellers a competitive edge.\n"
    "2. Statistical data features allow sellers to make informed decisions.\n"
    "3. The user-friendly interface ensures accessibility for non-technical sellers.\n"
    "4. Future iterations will focus on improving personalized insights for sellers.\n"
)

# Save the presentation
prs.save('ONDC_Application_Development_Plan.pptx')
